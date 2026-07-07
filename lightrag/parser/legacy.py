from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from lightrag.constants import PARSER_ENGINE_LEGACY
from lightrag.sidecar import write_sidecar
from lightrag.sidecar.ir import IRBlock, IRDoc, IRPosition
from lightrag.utils_pipeline import normalize_document_file_path, parsed_artifact_dir_for


_UTF8_TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".mdx",
    ".rtf",
    ".tex",
    ".html",
    ".htm",
    ".csv",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
    ".conf",
    ".ini",
    ".properties",
    ".sql",
    ".bat",
    ".sh",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".swift",
    ".go",
    ".rb",
    ".php",
    ".css",
    ".scss",
    ".less",
}


def extract_legacy_text(source_path: str | Path, *, password: str | None = None) -> str:
    """Extract text with the local legacy/direct-text parser.

    Text, data, and code files are decoded as UTF-8.  PDF/DOCX/PPTX/XLSX use
    the same lightweight local extraction strategy that the legacy API route
    historically used, without importing the router layer back into the
    pipeline.
    """

    path = Path(source_path)
    suffix = path.suffix.lower()
    if suffix in _UTF8_TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8")

    file_bytes = path.read_bytes()
    if suffix == ".pdf":
        return _extract_pdf(file_bytes, password=password)
    if suffix == ".docx":
        return _extract_docx(file_bytes)
    if suffix == ".pptx":
        return _extract_pptx(file_bytes)
    if suffix == ".xlsx":
        return _extract_xlsx(file_bytes)

    raise ValueError(f"Legacy parser does not support {suffix or 'extensionless'} files")


def parse_legacy_source_file(
    *,
    doc_id: str,
    file_path: str | Path,
    document_name: str | None = None,
    password: str | None = None,
) -> dict[str, Any]:
    """Extract source text and write a LightRAG sidecar next to the source file."""

    source_path = Path(file_path)
    if not source_path.is_file():
        raise FileNotFoundError(f"Document source not found: {source_path}")

    normalized_name = document_name or normalize_document_file_path(str(source_path))
    if normalized_name == "unknown_source":
        normalized_name = source_path.name or f"{doc_id}.txt"

    content = extract_legacy_text(source_path, password=password)
    parsed_dir = parsed_artifact_dir_for(normalized_name, parent_hint=source_path.parent)
    ir = IRDoc(
        document_name=normalized_name,
        document_format=source_path.suffix.lower().lstrip(".") or "text",
        doc_title=source_path.stem or normalized_name,
        split_option={"type": "legacy"},
        blocks=[
            IRBlock(
                content_template=content,
                positions=[
                    IRPosition(
                        type="absolute",
                        anchor=str(source_path.name),
                        range=[0, len(content)],
                    )
                ],
            )
        ],
    )
    parsed_data = write_sidecar(
        ir,
        parsed_dir=parsed_dir,
        doc_id=doc_id,
        engine=PARSER_ENGINE_LEGACY,
    )
    parsed_data["file_path"] = str(source_path)
    parsed_data["parse_engine"] = PARSER_ENGINE_LEGACY
    parsed_data["parse_stage_skipped"] = False
    return parsed_data


def _extract_pdf(file_bytes: bytes, *, password: str | None = None) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(BytesIO(file_bytes))
    if reader.is_encrypted:
        decrypt_result = reader.decrypt(password or "")
        if decrypt_result == 0:
            if password:
                raise ValueError("Incorrect PDF password")
            raise ValueError("PDF is encrypted but no password provided")

    content = ""
    for page in reader.pages:
        content += (page.extract_text() or "") + "\n"
    return content


def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document  # type: ignore
    from docx.table import Table  # type: ignore
    from docx.text.paragraph import Paragraph  # type: ignore

    doc = Document(BytesIO(file_bytes))

    def escape_cell(cell_value: str | None) -> str:
        if cell_value is None:
            return ""
        text = str(cell_value)
        return (
            text
            .replace("\\", "\\\\")
            .replace("\t", "&emsp;&emsp;")
            .replace("\r\n", "<br>")
            .replace("\r", "<br>")
            .replace("\n", "<br>")
        )

    content_parts: list[str] = []
    in_table = False
    for element in doc.element.body:
        if element.tag.endswith("p"):
            if in_table:
                content_parts.append("")
                in_table = False
            content_parts.append(Paragraph(element, doc).text)
        elif element.tag.endswith("tbl"):
            if content_parts and not in_table:
                content_parts.append("")
            in_table = True
            table = Table(element, doc)
            for row in table.rows:
                row_text = [escape_cell(cell.text) for cell in row.cells]
                if any(cell for cell in row_text):
                    content_parts.append("\t".join(row_text))
    return "\n".join(content_parts)


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(BytesIO(file_bytes))
    content_parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            shape_text = getattr(shape, "text", None)
            if isinstance(shape_text, str):
                content_parts.append(shape_text)
    return "\n".join(content_parts)


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook  # type: ignore

    workbook = load_workbook(BytesIO(file_bytes))

    def escape_cell(cell_value: Any) -> str:
        if cell_value is None:
            return ""
        text = str(cell_value)
        return (
            text
            .replace("\\", "\\\\")
            .replace("\t", "\\t")
            .replace("\r\n", "\\n")
            .replace("\r", "\\n")
            .replace("\n", "\\n")
        )

    content_parts: list[str] = []
    sheet_separator = "=" * 20
    for index, sheet in enumerate(workbook):
        if index > 0:
            content_parts.append("")
        safe_title = (
            str(sheet.title).replace("\n", " ").replace("\t", " ").replace("\r", " ")
        )
        content_parts.append(f"{sheet_separator} Sheet: {safe_title} {sheet_separator}")
        max_columns = sheet.max_column if sheet.max_column else 0
        for row in sheet.iter_rows(values_only=True):
            row_parts = [
                escape_cell(row[i]) if i < len(row) else ""
                for i in range(max_columns)
            ]
            content_parts.append(
                "" if all(part == "" for part in row_parts) else "\t".join(row_parts)
            )
    content_parts.append(sheet_separator)
    return "\n".join(content_parts)
